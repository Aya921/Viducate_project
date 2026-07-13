from fastapi import HTTPException,status
import fitz  
from pptx import Presentation

def extract_from_pdf(file_path: str) -> list[str]:
    slides_text = []
    doc = fitz.open(file_path)
    
    for page in doc:
        text = page.get_text().strip()
        if text:  
            slides_text.append(text)
    
    doc.close()
    return slides_text


def extract_from_pptx(file_path: str) -> list[str]:
    slides_text = []
    prs = Presentation(file_path)
    
    for slide in prs.slides:
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        
        if texts: 
            slides_text.append("\n".join(texts))
    
    return slides_text


def extract_slides_text(file_path: str) -> list[str]:
    if file_path.endswith(".pdf"):
        return extract_from_pdf(file_path)
    elif file_path.endswith(".pptx"):
        return extract_from_pptx(file_path)
    else:
        raise HTTPException(
            status_code=status.HTTP_400_BAD_REQUEST,
            detail=f"Unsupported format. Allowed: .pdf, .pptx"
        )